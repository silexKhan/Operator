#
#  auditor.py - Professional PRD & Planning Document Auditor
#

import re
import os
from mcp_operator.engine.interfaces import BaseAuditor

class MarkdownAuditor(BaseAuditor):
    def __init__(self, logger=None, circuit_manager=None):
        super().__init__(logger)
        self.circuit_manager = circuit_manager

    def _extract_text_from_pptx(self, file_path: str) -> str:
        """PPTX 파일에서 모든 슬라이드의 텍스트를 추출하여 마크다운 형식으로 변환합니다."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            full_text = []
            
            for i, slide in enumerate(prs.slides):
                full_text.append(f"## Slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 제목 텍스트 박스일 경우 헤더 처리
                        if shape == slide.shapes[0]: 
                            full_text.append(f"# {shape.text}")
                        else:
                            full_text.append(shape.text)
            
            return "\n".join(full_text)
        except Exception as e:
            self.log(f" PPTX 추출 실패: {str(e)}", 2)
            return ""

    def _get_protocols_for_path(self, file_path: str):
        """파일 경로 기반 기획 프로토콜 로드"""
        if self.circuit_manager:
            active = self.circuit_manager.get_active_circuit()
            if active:
                return active.get_protocols()
        return None

    def audit(self, file_path: str, content: str = "") -> list[str]:
        results = []
        
        # 0. PPTX 대응: 파일 확장자가 .pptx인 경우 텍스트 직접 추출
        if file_path.lower().endswith(".pptx"):
            self.log(f" PPTX 회선 연결됨. 텍스트 추출을 시작합니다... ({os.path.basename(file_path)})", 1)
            content = self._extract_text_from_pptx(file_path)
            if not content:
                return [" FAIL: PPTX 파일에서 텍스트를 추출할 수 없거나 파일이 손상되었습니다."]
        
        protocol_class = self._get_protocols_for_path(file_path)
        
        # 1. 필수 섹션 검사 (Protocol: PRD Integrity)
        # PPTX의 경우 슬라이드 제목이나 본문에 해당 키워드가 있는지 검사
        required_headers = ["목적", "배경", "상세 기능", "예외 케이스", "히스토리"]
        for header in required_headers:
            if header not in content: # PPTX는 정규식 대신 포함 여부로 유연하게 체크
                results.append(f" FAIL: Protocol (PRD Integrity) - 필수 항목 [{header}]가 누락되었습니다. ")

        # 2. 용어 통일성 검사 (Protocol: Terminology)
        # 예: '유저' 대신 '사용자' 권장
        if "유저" in content:
            results.append(" WARNING: Protocol (Terminology) - '유저' 대신 '사용자'라는 용어를 사용하십시오. ")

        # 3. Global Protocol 0-1 (Mechanical Integrity) 체크
        if "..." in content or "중략" in content:
            results.append(" FAIL: Protocol 0-1 (Global) - 문서 내에 '...' 또는 '중략'을 사용하지 마십시오. 완전한 명세가 필요합니다. ")

        # 4. 일자 형식 검사 (Format)
        if re.search(r"\d{2}\.\d{2}\.\d{2}", content):
            results.append(" [Format Reminder] 날짜는 'YYYY-MM-DD' 형식을 권장합니다.")

        return results
